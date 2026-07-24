"""Ch14-04: Export to PPTX - build a briefing deck with python-pptx (in-memory PNG).

Bilingual scene (zh/en) driven by ``EpiBaseScene.t()``. All on-screen prose is
read from ``TEXT`` via ``self.t(key)``; code strings stay identical (ASCII) across
languages. The generated slide is drawn as a mockup from ``RoundedRectangle`` /
``Rectangle`` / ``Line`` / ``Text`` (no new mobjects).
"""

from __future__ import annotations

from manim import (
    DOWN,
    LEFT,
    ORIGIN,
    RIGHT,
    UP,
    FadeIn,
    FadeOut,
    Line,
    ManimColor,
    Rectangle,
    RoundedRectangle,
    Text,
    VGroup,
)

from videos.src.base_scene import EpiBaseScene
from videos.src.code_mobjects import (
    ACCENT_BLUE,
    ACCENT_ORANGE,
    BG_CARD,
    BG_CARD_ALT,
    BORDER_LIGHT,
    FONT_CJK,
    TEXT_PRIMARY,
    TEXT_SECONDARY,
    BlindSpotBanner,
    ExtraExampleBanner,
)


class Ch14ExportPptxScene(EpiBaseScene):
    """Tutorial video scene: exporting a PPTX briefing deck with python-pptx."""

    total_steps: int = 9

    TEXT: dict[str, dict[str, str]] = {
        "zh": {
            "title_main": "一鍵輸出 PPTX 簡報",
            "title_sub": "python-pptx：記憶體 PNG → 投影片",
            "why_heading": "為什麼要程式化輸出？",
            "why_lines": [
                "手動截圖貼投影片：慢、易貼錯版本",
                "資料一更新 → 整套重做",
                "python-pptx 直接組簡報，零複製貼上",
                "→ 可重現：重跑 notebook，簡報重生",
            ],
            "bytesio_heading": "第一步：圖存進記憶體 PNG",
            "bytesio_title": "export_pptx.py",
            "pptx_code_heading": "python-pptx：一行貼上圖表",
            "pptx_code_title": "export_pptx.py",
            "slide_mockup_heading": "產出的簡報長這樣",
            "slide_title": "退伍軍人症群聚 SitRep",
            "slide_pic_caption": "圖：淋浴 vs 侵襲率",
            "slide_table_rows": ["住民　　280", "感染　　43.2%", "死亡　　15.7%"],
            "summary_heading": "PPTX 輸出三重點",
            "summary_lines": [
                "① io.BytesIO：圖表不落地就能塞進簡報",
                "② add_slide + add_picture 是核心兩招",
                "③ 程式化 = 一鍵重跑、格式一致、零手工錯",
                "→ 資料更新，簡報自動同步",
            ],
            "extra_banner_title": "額外範例：自動產生的流感週報簡報",
            "extra_auto_heading": "包成函式，自動生整份簡報",
            "extra_auto_title": "auto_sitrep.py",
            "blindspot_banner_title": "PPTX 輸出三個新手地雷",
            "outro_heading": "下一集：輸出正式 DOCX 報告",
            "outro_sub": "python-docx：標題、表格、內嵌圖",
        },
        "en": {
            "title_main": "Export a PPTX Deck in One Click",
            "title_sub": "python-pptx: in-memory PNG -> slide",
            "why_heading": "Why Export Programmatically?",
            "why_lines": [
                "Manual screenshots: slow, wrong-version prone",
                "Data changes -> redo the whole deck",
                "python-pptx builds the deck, zero copy-paste",
                "-> reproducible: rerun the notebook, deck rebuilds",
            ],
            "bytesio_heading": "Step 1: save the figure to an in-memory PNG",
            "bytesio_title": "export_pptx.py",
            "pptx_code_heading": "python-pptx: drop the chart in one line",
            "pptx_code_title": "export_pptx.py",
            "slide_mockup_heading": "What the Deck Looks Like",
            "slide_title": "Legionnaires' cluster SitRep",
            "slide_pic_caption": "Fig: shower vs attack rate",
            "slide_table_rows": ["Residents  280", "Infected   43.2%", "Deaths     15.7%"],
            "summary_heading": "Three Takeaways on PPTX Export",
            "summary_lines": [
                "1. io.BytesIO: chart to deck without hitting disk",
                "2. add_slide + add_picture are the two core moves",
                "3. Programmatic = one rerun, consistent, no hand errors",
                "-> data updates, the deck syncs itself",
            ],
            "extra_banner_title": "Extra example: an auto-generated flu SitRep deck",
            "extra_auto_heading": "Wrap it in a function, auto-build the deck",
            "extra_auto_title": "auto_sitrep.py",
            "blindspot_banner_title": "Three Beginner Blind Spots",
            "outro_heading": "Next: export a formal DOCX report",
            "outro_sub": "python-docx: headings, table, embedded figure",
        },
    }

    def construct(self) -> None:
        self.construct_from_segments()

    # ------------------------------------------------------------------
    # Shared helpers
    # ------------------------------------------------------------------

    def _bullets(self, heading_key: str, lines_key: str, duration: float) -> None:
        heading = self.t(heading_key)
        lines = self.t(lines_key)
        h = Text(heading, font=FONT_CJK, font_size=30, color=ACCENT_ORANGE).to_edge(UP, buff=0.8)
        bl = (
            VGroup(*[Text(x, font=FONT_CJK, font_size=22, color=TEXT_PRIMARY) for x in lines])
            .arrange(DOWN, aligned_edge=LEFT, buff=0.4)
            .next_to(h, DOWN, buff=0.55)
        )
        if bl.width > 12.5:
            bl.scale_to_fit_width(12.5)
        self.play(FadeIn(h), run_time=0.5)
        self.play(FadeIn(bl, lag_ratio=0.2), run_time=1.2)
        self.wait(max(0.1, duration - 1.7))
        self.play(FadeOut(VGroup(h, bl)), run_time=0.5)

    def _code_block(self, heading_key: str, title_key: str, code: str, duration: float) -> None:
        h = Text(self.t(heading_key), font=FONT_CJK, font_size=26, color=ACCENT_ORANGE).to_edge(
            UP, buff=0.5
        )
        self.play(FadeIn(h), run_time=0.4)
        panel = self.show_code(code, title=self.t(title_key), position=ORIGIN + DOWN * 0.3)
        self.wait(max(0.1, duration - 1.4))
        self.play(FadeOut(VGroup(h, panel)), run_time=0.5)

    # ------------------------------------------------------------------
    # Main lesson
    # ------------------------------------------------------------------

    def show_title(self, duration: float = 3.0, **kwargs) -> None:
        self.show_title_card(self.t("title_main"), self.t("title_sub"), duration=duration)

    def show_why_export(self, duration: float = 7.0, **kwargs) -> None:
        self.show_step_indicator(1, self.total_steps)
        self._bullets("why_heading", "why_lines", duration)

    def show_bytesio(self, duration: float = 8.0, **kwargs) -> None:
        self.show_step_indicator(2, self.total_steps)
        code = kwargs.get(
            "code",
            (
                "import io\n"
                "buf = io.BytesIO()\n"
                'fig.savefig(buf, format="png", dpi=150)\n'
                "plt.close(fig)\n"
                "buf.seek(0)   # rewind before reading"
            ),
        )
        self._code_block("bytesio_heading", "bytesio_title", code, duration)

    def show_pptx_code(self, duration: float = 9.0, **kwargs) -> None:
        self.show_step_indicator(3, self.total_steps)
        code = kwargs.get(
            "code",
            (
                "from pptx import Presentation\n"
                "from pptx.util import Inches\n"
                "prs = Presentation()\n"
                "t = prs.slides.add_slide(prs.slide_layouts[0])\n"
                't.shapes.title.text = "Legionella cluster SitRep"\n'
                "s = prs.slides.add_slide(prs.slide_layouts[5])\n"
                "s.shapes.add_picture(buf, Inches(0.5), Inches(1.2))\n"
                'prs.save("legionella_sitrep.pptx")'
            ),
        )
        self._code_block("pptx_code_heading", "pptx_code_title", code, duration)

    def show_slide_mockup(self, duration: float = 8.0, **kwargs) -> None:
        self.show_step_indicator(4, self.total_steps)

        heading = Text(
            self.t("slide_mockup_heading"), font=FONT_CJK, font_size=28, color=ACCENT_ORANGE
        ).to_edge(UP, buff=0.5)

        slide = RoundedRectangle(
            corner_radius=0.12,
            width=9.4,
            height=5.0,
            fill_color=ManimColor(BG_CARD),
            fill_opacity=1,
            stroke_color=ManimColor(BORDER_LIGHT),
            stroke_width=2,
        )
        stitle = Text(
            self.t("slide_title"), font=FONT_CJK, font_size=22, color=TEXT_PRIMARY, weight="BOLD"
        )
        stitle.move_to(slide.get_top() + DOWN * 0.55)
        stitle.align_to(slide, LEFT).shift(RIGHT * 0.6)
        underline = Line(
            ORIGIN, RIGHT * 4.2, color=ManimColor(ACCENT_ORANGE), stroke_width=3
        ).next_to(stitle, DOWN, buff=0.12, aligned_edge=LEFT)

        # Picture box with two mini bars (left)
        picbox = Rectangle(
            width=4.2,
            height=2.6,
            fill_color=ManimColor(BG_CARD_ALT),
            fill_opacity=1,
            stroke_color=ManimColor(BORDER_LIGHT),
            stroke_width=1.5,
        )
        bar1 = Rectangle(
            width=0.8, height=1.0, fill_color=ManimColor(ACCENT_BLUE), fill_opacity=1, stroke_width=0
        )
        bar2 = Rectangle(
            width=0.8,
            height=1.8,
            fill_color=ManimColor(ACCENT_ORANGE),
            fill_opacity=1,
            stroke_width=0,
        )
        bars = VGroup(bar1, bar2).arrange(RIGHT, buff=0.6, aligned_edge=DOWN)
        bars.move_to(picbox.get_center())
        bars.align_to(picbox, DOWN).shift(UP * 0.35)
        piccap = Text(
            self.t("slide_pic_caption"), font=FONT_CJK, font_size=15, color=TEXT_SECONDARY
        ).next_to(picbox, DOWN, buff=0.12)
        left_group = VGroup(picbox, bars, piccap)

        # Mini summary table (right)
        rows = self.t("slide_table_rows")
        tbl_lines = VGroup(
            *[Text(r, font=FONT_CJK, font_size=17, color=TEXT_PRIMARY) for r in rows]
        ).arrange(DOWN, aligned_edge=LEFT, buff=0.28)
        tbl_card = RoundedRectangle(
            corner_radius=0.1,
            width=3.2,
            height=2.4,
            fill_color=ManimColor(BG_CARD_ALT),
            fill_opacity=1,
            stroke_color=ManimColor(BORDER_LIGHT),
            stroke_width=1.5,
        )
        if tbl_lines.width > 2.8:
            tbl_lines.scale_to_fit_width(2.8)
        tbl_lines.move_to(tbl_card.get_center())
        right_group = VGroup(tbl_card, tbl_lines)

        body = VGroup(left_group, right_group).arrange(RIGHT, buff=0.5)
        body.move_to(slide.get_center() + DOWN * 0.35)

        full = VGroup(slide, stitle, underline, body)

        self.play(FadeIn(heading), run_time=0.4)
        self.play(FadeIn(slide), FadeIn(stitle), FadeIn(underline), run_time=0.7)
        self.play(FadeIn(body, lag_ratio=0.1), run_time=1.1)
        self.wait(max(0.1, duration - 2.7))
        self.play(FadeOut(VGroup(heading, full)), run_time=0.5)

    def show_main_summary(self, duration: float = 6.0, **kwargs) -> None:
        self.show_step_indicator(5, self.total_steps)
        self._bullets("summary_heading", "summary_lines", duration)

    # ------------------------------------------------------------------
    # Extra example
    # ------------------------------------------------------------------

    def show_extra_banner(self, duration: float = 2.0, **kwargs) -> None:
        self.show_section_banner(
            ExtraExampleBanner(self.t("extra_banner_title")), duration=duration
        )

    def show_extra_auto_deck(self, duration: float = 8.0, **kwargs) -> None:
        self.show_step_indicator(6, self.total_steps)
        code = kwargs.get(
            "code",
            (
                "def build_sitrep(figs):\n"
                "    prs = Presentation()\n"
                "    for fig in figs:\n"
                "        s = prs.slides.add_slide(prs.slide_layouts[5])\n"
                "        s.shapes.add_picture(to_png(fig), Inches(0.5))\n"
                '    prs.save("weekly_flu_sitrep.pptx")'
            ),
        )
        self._code_block("extra_auto_heading", "extra_auto_title", code, duration)

    # ------------------------------------------------------------------
    # Blind spots
    # ------------------------------------------------------------------

    def show_blindspot_banner(self, duration: float = 2.0, **kwargs) -> None:
        self.show_section_banner(
            BlindSpotBanner(self.t("blindspot_banner_title")), duration=duration
        )

    def show_blindspot_forgot_seek(self, duration: float = 6.0, **kwargs) -> None:
        self.show_step_indicator(7, self.total_steps)
        panel = self.show_error_vs_correct(
            kwargs.get("error_code", "slide.shapes.add_picture(buf)"),
            kwargs.get("correct_code", "buf.seek(0); slide.shapes.add_picture(buf)"),
            duration=duration,
        )
        self.play(FadeOut(panel), run_time=0.5)

    def show_blindspot_manual_screenshot(self, duration: float = 6.0, **kwargs) -> None:
        self.show_step_indicator(8, self.total_steps)
        panel = self.show_error_vs_correct(
            kwargs.get("error_code", 'img = screenshot("chart.png")'),
            kwargs.get("correct_code", 'fig.savefig(buf, format="png")'),
            duration=duration,
        )
        self.play(FadeOut(panel), run_time=0.5)

    def show_blindspot_forgot_close(self, duration: float = 6.0, **kwargs) -> None:
        self.show_step_indicator(9, self.total_steps)
        panel = self.show_error_vs_correct(
            kwargs.get("error_code", "fig.savefig(buf)"),
            kwargs.get("correct_code", "fig.savefig(buf); plt.close(fig)"),
            duration=duration,
        )
        self.play(FadeOut(panel), run_time=0.5)

    # ------------------------------------------------------------------
    # Outro
    # ------------------------------------------------------------------

    def show_outro(self, duration: float = 3.0, **kwargs) -> None:
        self.show_step_indicator(self.total_steps, self.total_steps)
        h = Text(self.t("outro_heading"), font=FONT_CJK, font_size=26, color=ACCENT_ORANGE).move_to(
            ORIGIN + UP * 0.5
        )
        s = Text(self.t("outro_sub"), font=FONT_CJK, font_size=20, color=TEXT_SECONDARY).next_to(
            h, DOWN, buff=0.4
        )
        self.play(FadeIn(h), run_time=0.6)
        self.play(FadeIn(s), run_time=0.5)
        self.wait(max(0.1, duration - 1.1))
        self.play(FadeOut(VGroup(h, s)), run_time=0.5)
