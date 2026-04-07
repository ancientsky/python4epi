"""Ch04-07: Report Output — 專業報告輸出

Manim scene for the tutorial video on generating professional reports
(Dashboard, Word, PPT, PDF) from SitRep data.
"""

from __future__ import annotations

from manim import (
    DOWN,
    LEFT,
    RIGHT,
    UP,
    ORIGIN,
    FadeIn,
    FadeOut,
    Text,
    VGroup,
)

from videos.src.base_scene import EpiBaseScene
from videos.src.code_mobjects import (
    ACCENT_ORANGE,
    FONT_CJK,
    FONT_MONO,
    TEXT_PRIMARY,
    TEXT_SECONDARY,
    BlindSpotBanner,
    ExtraExampleBanner,
)


class Ch04ReportOutputScene(EpiBaseScene):
    """Tutorial video scene: professional report output formats."""

    total_steps: int = 16

    def construct(self) -> None:
        self.construct_from_segments()

    # ------------------------------------------------------------------
    # Main lesson
    # ------------------------------------------------------------------

    def show_title(self, duration: float = 3.0, **kwargs) -> None:
        """Title card for the Report Output lesson."""
        self.show_title_card("專業報告輸出", "Dashboard、Word、PPT、PDF 一次搞定", duration=duration)

    def show_bytesio_concept(self, duration: float = 6.0, **kwargs) -> None:
        """Show BytesIO for in-memory file handling."""
        self.show_step_indicator(1, self.total_steps)

        code_lines = kwargs.get(
            "code",
            (
                "from io import BytesIO\n"
                "\n"
                "buf = BytesIO()\n"
                "fig.savefig(buf, format='png', dpi=150)\n"
                "buf.seek(0)  # rewind to start\n"
                "# buf can now be embedded in Word/PPT"
            ),
        )

        self.show_code(code_lines, title="bytesio.py")
        self.wait(duration)
        self.clear_screen()

    def show_plotly_dashboard(self, duration: float = 6.0, **kwargs) -> None:
        """Show Plotly dashboard creation."""
        self.show_step_indicator(2, self.total_steps)

        code_lines = kwargs.get(
            "code",
            (
                "import plotly.graph_objects as go\n"
                "from plotly.subplots import make_subplots\n"
                "\n"
                "fig = make_subplots(rows=2, cols=2,\n"
                "    subplot_titles=('Epi Curve','Age Dist',\n"
                "                    'Wing AR','KPI'))\n"
                "fig.update_layout(title='SitRep Dashboard')\n"
                "fig.write_html('dashboard.html')"
            ),
        )

        self.show_code(code_lines, title="dashboard.py")
        self.wait(duration)
        self.clear_screen()

    def show_docx_report(self, duration: float = 6.0, **kwargs) -> None:
        """Show python-docx Word report generation."""
        self.show_step_indicator(3, self.total_steps)

        code_lines = kwargs.get(
            "code",
            (
                "from docx import Document\n"
                "\n"
                "doc = Document()\n"
                "doc.add_heading('SitRep 2026-01-20', level=1)\n"
                "doc.add_paragraph(\n"
                "    f'Attack Rate: {ar_pct}%'\n"
                ")\n"
                "doc.add_picture(buf, width=Inches(5))\n"
                "doc.save('sitrep.docx')"
            ),
        )

        self.show_code(code_lines, title="docx_report.py")
        self.wait(duration)
        self.clear_screen()

    def show_pptx_slides(self, duration: float = 6.0, **kwargs) -> None:
        """Show python-pptx slide generation."""
        self.show_step_indicator(4, self.total_steps)

        code_lines = kwargs.get(
            "code",
            (
                "from pptx import Presentation\n"
                "from pptx.util import Inches\n"
                "\n"
                "prs = Presentation()\n"
                "slide = prs.slides.add_slide(\n"
                "    prs.slide_layouts[5]\n"
                ")\n"
                "slide.shapes.add_picture(\n"
                "    buf, Inches(1), Inches(1.5),\n"
                "    width=Inches(8)\n"
                ")\n"
                "prs.save('sitrep.pptx')"
            ),
        )

        self.show_code(code_lines, title="pptx_slides.py")
        self.wait(duration)
        self.clear_screen()

    def show_pdf_report(self, duration: float = 6.0, **kwargs) -> None:
        """Show PDF generation with matplotlib."""
        self.show_step_indicator(5, self.total_steps)

        code_lines = kwargs.get(
            "code",
            (
                "from matplotlib.backends.backend_pdf import PdfPages\n"
                "\n"
                "with PdfPages('sitrep.pdf') as pdf:\n"
                "    # page 1: epi curve\n"
                "    pdf.savefig(fig_curve)\n"
                "    # page 2: age distribution\n"
                "    pdf.savefig(fig_age)\n"
                "    # page 3: wing attack rates\n"
                "    pdf.savefig(fig_wing)"
            ),
        )

        self.show_code(code_lines, title="pdf_report.py")
        self.wait(duration)
        self.clear_screen()

    def show_main_summary(self, duration: float = 5.0, **kwargs) -> None:
        """Summarise key points about report output."""
        self.show_step_indicator(6, self.total_steps)

        heading = Text(
            "重點整理",
            font=FONT_CJK,
            font_size=34,
            color=ACCENT_ORANGE,
        ).to_edge(UP, buff=0.8)

        points = VGroup(
            Text("1. BytesIO 讓圖片在記憶體中傳遞，不需存檔", font=FONT_CJK, font_size=23, color=TEXT_PRIMARY),
            Text("2. Plotly Dashboard 適合互動式報告", font=FONT_CJK, font_size=23, color=TEXT_PRIMARY),
            Text("3. python-docx 產生 Word 報告", font=FONT_CJK, font_size=23, color=TEXT_PRIMARY),
            Text("4. python-pptx 產生簡報投影片", font=FONT_CJK, font_size=23, color=TEXT_PRIMARY),
            Text("5. PdfPages 可將多張圖合併成一份 PDF", font=FONT_CJK, font_size=23, color=TEXT_PRIMARY),
        ).arrange(DOWN, aligned_edge=LEFT, buff=0.40).next_to(heading, DOWN, buff=0.6)

        self.play(FadeIn(heading), run_time=0.5)
        self.play(FadeIn(points, lag_ratio=0.25), run_time=1.2)
        self.wait(duration - 1.7)
        self.play(FadeOut(VGroup(heading, points)), run_time=0.5)

    # ------------------------------------------------------------------
    # Extra example
    # ------------------------------------------------------------------

    def show_extra_banner(self, duration: float = 2.0, **kwargs) -> None:
        """Show the ExtraExampleBanner section divider."""
        banner = ExtraExampleBanner("額外範例：自動寄信附報告")
        self.show_section_banner(banner, duration=duration)

    def show_extra_email(self, duration: float = 6.0, **kwargs) -> None:
        """Email automation example."""
        self.show_step_indicator(8, self.total_steps)

        code_lines = kwargs.get(
            "code",
            (
                "import smtplib\n"
                "from email.message import EmailMessage\n"
                "\n"
                "msg = EmailMessage()\n"
                "msg['Subject'] = 'SitRep 2026-01-20'\n"
                "msg['To'] = 'epi-team@example.com'\n"
                "with open('sitrep.pdf', 'rb') as f:\n"
                "    msg.add_attachment(\n"
                "        f.read(), maintype='application',\n"
                "        subtype='pdf', filename='sitrep.pdf'\n"
                "    )"
            ),
        )

        self.show_code(code_lines, title="email_report.py")
        self.wait(duration)
        self.clear_screen()

    # ------------------------------------------------------------------
    # Blind spots
    # ------------------------------------------------------------------

    def show_blindspot_banner(self, duration: float = 2.0, **kwargs) -> None:
        """Show the BlindSpotBanner section divider."""
        banner = BlindSpotBanner("初學者常見地雷 3 選")
        self.show_section_banner(banner, duration=duration)

    def show_blindspot_seek(self, duration: float = 5.0, **kwargs) -> None:
        """Error vs correct: forgetting buf.seek(0) after writing."""
        error_code = kwargs.get("error_code", "fig.savefig(buf); doc.add_picture(buf)")
        correct_code = kwargs.get("correct_code", "fig.savefig(buf); buf.seek(0); doc.add_picture(buf)")
        panel = self.show_error_vs_correct(error_code, correct_code, duration=duration)
        self.play(FadeOut(panel), run_time=0.5)

    def show_blindspot_import(self, duration: float = 5.0, **kwargs) -> None:
        """Error vs correct: wrong import for Inches."""
        error_code = kwargs.get("error_code", "from pptx import Inches  # ImportError")
        correct_code = kwargs.get("correct_code", "from pptx.util import Inches  # correct")
        panel = self.show_error_vs_correct(error_code, correct_code, duration=duration)
        self.play(FadeOut(panel), run_time=0.5)

    def show_blindspot_cjk(self, duration: float = 5.0, **kwargs) -> None:
        """Error vs correct: CJK font not embedded in PDF."""
        error_code = kwargs.get("error_code", "plt.savefig('report.pdf')  # tofu boxes")
        correct_code = kwargs.get("correct_code", "plt.rcParams['font.sans-serif'] = ['Noto Sans CJK']")
        panel = self.show_error_vs_correct(error_code, correct_code, duration=duration)
        self.play(FadeOut(panel), run_time=0.5)

    # ------------------------------------------------------------------
    # Outro
    # ------------------------------------------------------------------

    def show_outro(self, duration: float = 3.0, **kwargs) -> None:
        """Closing card."""
        self.show_step_indicator(self.total_steps, self.total_steps)

        heading = Text(
            "Chapter 04 完結！",
            font=FONT_CJK,
            font_size=28,
            color=ACCENT_ORANGE,
        ).move_to(ORIGIN + UP * 0.5)

        sub = Text(
            "你已經學會從 Line List 產出完整 SitRep！",
            font=FONT_CJK,
            font_size=22,
            color=TEXT_SECONDARY,
        ).next_to(heading, DOWN, buff=0.4)

        self.play(FadeIn(heading), run_time=0.6)
        self.play(FadeIn(sub), run_time=0.5)
        self.wait(duration - 1.1)
        self.play(FadeOut(VGroup(heading, sub)), run_time=0.5)
